import os
import sys
from pptx import Presentation
import json
import subprocess
from pdf2image import convert_from_path

def convert_pptx_to_pdf(pptx_path, pdf_path):
    try:
        output_dir = os.path.dirname(pdf_path)
        subprocess.run(['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', 
                        output_dir, pptx_path], 
                       check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
        
        # Check if the PDF was actually created
        expected_pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + '.pdf'
        expected_pdf_path = os.path.join(output_dir, expected_pdf_name)
        
        if not os.path.exists(expected_pdf_path):
            raise FileNotFoundError(f"PDF file was not created at {expected_pdf_path}")
        
        # If the PDF was created with a different name, rename it
        if expected_pdf_path != pdf_path:
            os.rename(expected_pdf_path, pdf_path)
        
        print(f"Converted {pptx_path} to PDF: {pdf_path}")
        return pdf_path
    except subprocess.CalledProcessError as e:
        print(f"Error converting to PDF: {e}")
        print(f"STDOUT: {e.stdout.decode()}")
        print(f"STDERR: {e.stderr.decode()}")
        raise
    except Exception as e:
        print(f"Unexpected error during PDF conversion: {e}")
        raise

def extract_slides_and_create_json(ppt_file):
    output_folder = "extracted_slides"
    os.makedirs(output_folder, exist_ok=True)

    pdf_path = os.path.join(output_folder, "presentation.pdf")
    
    try:
        # Convert PPTX to PDF
        pdf_path = convert_pptx_to_pdf(ppt_file, pdf_path)

        # Check if PDF exists
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF file not found at {pdf_path}")

        # Convert PDF to images
        images = convert_from_path(pdf_path)

        # Process slides and create JSON
        prs = Presentation(ppt_file)
        slides_data = []

        for index, (slide, image) in enumerate(zip(prs.slides, images), start=1):
            print(f"Processing slide {index}...")
            
            image_filename = f"slide_{index}.jpg"
            image_path = os.path.join(output_folder, image_filename)
            image.save(image_path, "JPEG")

            # Get the speaker notes
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text if notes_slide else ""

            # Append slide data to the list
            slides_data.append({
                "media": image_filename,
                "text": notes_text
            })

        # Create a single data.json file with all slide data
        json_path = os.path.join(output_folder, "data.json")
        with open(json_path, 'w', encoding='utf-8') as json_file:
            json.dump(slides_data, json_file, indent=4, ensure_ascii=False)

        print(f"Successfully extracted {len(prs.slides)} slides and created a single data.json file in the '{output_folder}' directory.")

    except Exception as e:
        print(f"An error occurred: {str(e)}")
        print("Error details:")
        import traceback
        traceback.print_exc()
    finally:
        # Clean up the temporary PDF file
        if os.path.exists(pdf_path):
            os.remove(pdf_path)

# Usage
if __name__ == "__main__":
    if len(sys.argv) > 1:
        ppt_file = sys.argv[1]
    else:
        ppt_file = "test.pptx"
    
    if not os.path.exists(ppt_file):
        print(f"Error: File '{ppt_file}' not found.")
        sys.exit(1)
    
    extract_slides_and_create_json(ppt_file)
